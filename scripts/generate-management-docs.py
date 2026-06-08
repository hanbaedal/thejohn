# -*- coding: utf-8 -*-
"""관리용 시스템 구조 Word·PowerPoint 생성."""
from __future__ import annotations

import os
from datetime import date

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(ROOT, "docs")
TODAY = date.today().strftime("%Y-%m-%d")


def add_heading(doc, text, level=1):
    doc.add_heading(text, level=level)


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr[i].text = h
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            table.rows[ri + 1].cells[ci].text = str(val)
    doc.add_paragraph()


def build_word(path):
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("더존(thejohn) 업무 시스템\n구조 설명서")
    run.bold = True
    run.font.size = Pt(22)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run(f"작성 기준일: {TODAY}\n관리·기획용 요약 문서").font.size = Pt(11)
    doc.add_page_break()

    add_heading(doc, "1. 문서 목적", 1)
    doc.add_paragraph(
        "본 문서는 개발 소스 상세가 아니라, 경영·관리 관점에서 "
        "「더존 홈페이지·업무 시스템」이 어떤 역할을 하고, "
        "누가 무엇을 사용하며, 데이터와 서비스가 어떻게 연결되는지를 설명합니다."
    )
    doc.add_paragraph(
        "기술 담당자용 상세는 저장소 내 docs/ARCHITECTURE.md 를 참고하세요."
    )

    add_heading(doc, "2. 시스템 개요", 1)
    doc.add_paragraph(
        "thejohn(더존)은 거래처(업체) 주문, 내부 직원의 상품·업체 관리, "
        "발주서·거래명세서 발행, 고객지원(소식·문의) 등을 하나의 웹에서 처리하는 통합 업무 플랫폼입니다."
    )
    add_table(
        doc,
        ["항목", "내용"],
        [
            ["서비스 형태", "웹 브라우저 접속 (PC·모바일)"],
            ["운영 URL", "thejohn.co.kr / www.thejohn.co.kr"],
            ["호스팅", "Render(클라우드) — 애플리케이션 서버"],
            ["데이터 저장", "MongoDB Atlas(클라우드 DB)"],
            ["소스·배포", "GitHub 저장소 → Render 자동 배포"],
        ],
    )

    add_heading(doc, "3. 이용자·권한 구조", 1)
    add_table(
        doc,
        ["역할", "주요 이용자", "할 수 있는 일 (요약)"],
        [
            [
                "슈퍼바이저",
                "본사·총괄 관리",
                "전체 조회, 관리자 계정 생성, 통계, 모든 주문·PDF·수기 거래명세",
            ],
            [
                "관리자(admin)",
                "영업·담당 직원",
                "담당 거래처·상품 관리, 주문 처리, 발주서·거래명세, 수기 명세(권한 시)",
            ],
            [
                "거래처(vendor)",
                "업체 로그인",
                "담당 관리자 상품 주문, 장바구니, 주문 이력·PDF 저장",
            ],
            [
                "게스트",
                "비로그인·열람",
                "상품 목록 열람(가격 제한), 접속 통계",
            ],
        ],
    )

    add_heading(doc, "4. 업무 기능 구성 (관리자 메뉴 기준)", 1)
    add_table(
        doc,
        ["업무 영역", "화면·기능 예", "비고"],
        [
            ["홈·회사 소개", "메인, 사업부문 소개", "공개"],
            ["상품·업체 마스터", "상품등록, 업체등록, 예비거래처, 엑셀 일괄", "MongoDB 저장"],
            ["거래처 주문", "상품 카탈로그, 장바구니, 주문", "업체 전용"],
            ["주문서 관리", "발주 목록, 거래명세 목록, PDF 보기·저장", "슈퍼바이저·관리자"],
            ["수기 거래명세", "작성·목록·PDF", "공급자·품목·발행일 관리"],
            ["직원·권한", "관리자 등록·리스트, 주문 권한 부여", "슈퍼바이저"],
            [
                "통계",
                "접속·이용, 디비사용(DB용량·이용시간), SOLAPI SMS 발송",
                "슈퍼바이저",
            ],
            ["문서", "매뉴얼·구조 설명서 다운로드", "슈퍼바이저"],
            ["고객지원", "소식, QnA, 1:1 문의, 자료실", "사업부문별"],
            ["메일", "업체 대상 메일 발송", "첨부파일 지원"],
        ],
    )

    add_heading(doc, "5. 시스템 계층 구조", 1)
    doc.add_paragraph("관리 관점의 3계층으로 보면 다음과 같습니다.")
    layers = (
        "【표현 계층】 웹 화면 (HTML·JavaScript)\n"
        "  → 사용자가 보는 목록, 입력 폼, PDF 보기 모달\n\n"
        "【업무·API 계층】 Node.js 서버 (Express)\n"
        "  → 로그인, 권한 검사, 주문·상품·업체 처리, PDF 생성\n\n"
        "【데이터 계층】 MongoDB\n"
        "  → 직원, 거래처, 상품, 주문, 수기 거래명세 등 영구 저장"
    )
    p = doc.add_paragraph(layers)
    p.paragraph_format.left_indent = Cm(0.5)

    add_heading(doc, "6. 주요 데이터(정보) 종류", 1)
    add_table(
        doc,
        ["데이터", "설명"],
        [
            ["직원(staff)", "내부 로그인 계정, 역할, 주문 권한, 인감·공급자 정보"],
            ["거래처(vendors)", "업체 계정, 등급·단가, 담당 관리자"],
            ["상품(products)", "품목, 가격, 사업부문, 사진 1장, 담당 관리자"],
            ["상품정보(product_info)", "식품 표시사항 등 항목별 입력값"],
            ["주문(orders)", "업체 주문·발주 내역, PDF 생성 여부"],
            ["수기 거래명세", "직원이 직접 작성한 거래명세 (DB 별도 컬렉션)"],
            ["고객지원", "소식, 게시글, 문의 내역"],
            ["접속·이용 로그", "로그인·페이지 방문·세션 종료 (access_logs)"],
            ["SOLAPI 발송 로그", "주문 SMS 발송 시도·성공·실패 (solapi_logs)"],
        ],
    )

    add_heading(doc, "6.1 업무관리(슈퍼바이저) 메뉴", 2)
    doc.add_paragraph("그룹 마케팅 관리 → 업무관리 화면 구성 (2026년 6월 기준):")
    add_table(
        doc,
        ["그룹", "메뉴", "설명"],
        [
            ["계정", "관리자 등록 / 관리자 리스트", "내부 직원 계정·권한"],
            ["통계", "접속·이용 통계 / 디비사용 통계", "이용 행태·DB 용량·건수 집계"],
            ["알림·문서", "SOLAPI 이용 현황 / 문서 다운로드", "SMS 발송 집계·오피스 문서"],
        ],
    )

    add_heading(doc, "7. 문서(PDF) 처리 방식", 1)
    doc.add_paragraph(
        "발주서·거래명세서는 서버에서 PDF를 생성합니다. "
        "「PDF 보기」는 화면 모달로 열람만 하고, "
        "「PDF 저장」을 선택할 때만 PC에 파일로 받습니다."
    )
    add_table(
        doc,
        ["문서", "생성 시점", "보기", "파일 저장"],
        [
            ["발주서", "주문 등록·조회 시", "모달", "「PDF 저장」 버튼"],
            ["주문 거래명세", "주문 기준", "모달", "「PDF 저장」 버튼"],
            ["수기 거래명세", "작성·저장 후", "모달", "작성 화면 「PDF 저장」"],
        ],
    )

    add_heading(doc, "8. 외부 연동·인프라", 1)
    add_table(
        doc,
        ["구분", "서비스", "용도"],
        [
            ["웹 호스팅", "Render", "프로그램 실행·HTTPS"],
            ["데이터베이스", "MongoDB Atlas", "모든 업무 데이터"],
            ["도메인·DNS", "가비아 등", "thejohn.co.kr 연결"],
            ["SMS(선택)", "SOLAPI", "업체 주문 시 담당 관리자 SMS. 이용 현황은 solapi_logs 집계"],
            ["메일(선택)", "SMTP", "업체 메일 발송"],
        ],
    )

    add_heading(doc, "9. 보안·운영 요약", 1)
    bullets = [
        "로그인 후 발급되는 토큰(JWT)으로 API 접근 — 역할별 메뉴·API 이중 제한",
        "DB·비밀번호·JWT는 Render 환경 변수로 관리 (소스에 비밀번호 미포함)",
        "서비스 상태 확인: /api/health (DB 연결 여부)",
        "배포: GitHub main 브랜치 반영 시 Render 자동 빌드·재시작",
    ]
    for b in bullets:
        doc.add_paragraph(b, style="List Bullet")

    add_heading(doc, "10. 문의·기술 참고", 1)
    doc.add_paragraph("저장소: github.com/hanbaedal/thejohn")
    doc.add_paragraph("운영 장애·DB: deploy/RENDER-FIX.md, deploy/MONGODB-FIX.md")
    doc.add_paragraph("개발 구조 상세: docs/ARCHITECTURE.md")

    doc.save(path)
    print("Word:", path)


def add_slide_title(prs, title, subtitle=None):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_slide_bullets(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = PptPt(18)


def add_slide_table(prs, title, headers, rows):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    rows_n = len(rows) + 1
    cols_n = len(headers)
    left, top, width, height = Inches(0.5), Inches(1.4), Inches(9), Inches(0.35 * rows_n)
    shape = slide.shapes.add_table(rows_n, cols_n, left, top, width, height)
    table = shape.table
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r + 1, c).text = str(val)


def build_ppt(path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_slide_title(
        prs,
        "더존(thejohn) 업무 시스템",
        f"시스템 구조 (관리용 요약)\n{TODAY}",
    )

    add_slide_bullets(
        prs,
        "1. 시스템이 하는 일",
        [
            "거래처 온라인 주문 · 내부 상품·업체 관리",
            "발주서·거래명세서 PDF 발행 (보기/저장 분리)",
            "수기 거래명세서 · 고객지원(소식·문의) · 통계",
            "웹(thejohn.co.kr) + 클라우드(DB·서버)",
        ],
    )

    add_slide_table(
        prs,
        "2. 이용자·권한",
        ["역할", "대상", "핵심 권한"],
        [
            ["슈퍼바이저", "총괄", "전체·계정·통계"],
            ["관리자", "담당 직원", "담당 거래처·주문·PDF"],
            ["거래처", "업체", "주문·이력"],
            ["게스트", "방문자", "상품 열람"],
        ],
    )

    add_slide_bullets(
        prs,
        "3. 업무 메뉴 구성",
        [
            "마스터: 상품(사진 1장·상품정보) · 업체 · 예비거래처",
            "주문: 거래처 장바구니 → 주문서 관리(발주·명세)",
            "문서: 발주서 PDF · 거래명세 PDF · 수기 명세",
            "지원: 소식 · QnA · 문의 · 메일",
        ],
    )

    add_slide_table(
        prs,
        "3-1. 업무관리(슈퍼바이저)",
        ["그룹", "메뉴"],
        [
            ["계정", "관리자 등록 · 관리자 리스트"],
            ["통계", "접속·이용 · 디비사용"],
            ["알림·문서", "SOLAPI 이용 현황 · 문서 다운로드"],
        ],
    )

    add_slide_bullets(
        prs,
        "4. 시스템 3계층",
        [
            "① 화면 — 브라우저(목록·입력·PDF 모달)",
            "② 서버 — Node.js(권한·업무·PDF 생성)",
            "③ DB — MongoDB(직원·거래처·상품·주문 등)",
        ],
    )

    add_slide_table(
        prs,
        "5. 인프라",
        ["항목", "내용"],
        [
            ["호스팅", "Render"],
            ["DB", "MongoDB Atlas"],
            ["배포", "GitHub → Render 자동"],
            ["도메인", "thejohn.co.kr"],
        ],
    )

    add_slide_bullets(
        prs,
        "6. PDF·보안 요약",
        [
            "PDF 보기: 화면 모달 (자동 다운로드 없음)",
            "PDF 저장: 사용자가 저장 선택 시만 파일 저장",
            "로그인·역할별 접근 제어 (JWT)",
            "상태 점검: /api/health",
        ],
    )

    add_slide_title(
        prs,
        "감사합니다",
        "상세: docs/ARCHITECTURE.md · docs/thejohn-system-structure-management.docx",
    )

    prs.save(path)
    print("PPT:", path)


def main():
    os.makedirs(DOCS, exist_ok=True)
    # ASCII 파일명 (Windows·Git 호환). 문서 제목은 한글.
    docx_path = os.path.join(DOCS, "thejohn-system-structure-management.docx")
    pptx_path = os.path.join(DOCS, "thejohn-system-structure-management.pptx")
    build_word(docx_path)
    build_ppt(pptx_path)


if __name__ == "__main__":
    main()
